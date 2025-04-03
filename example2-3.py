# basic import 
from mcp.server.fastmcp import FastMCP, Image
from mcp.server.fastmcp.prompts import base
from mcp.types import TextContent
from mcp import types
from PIL import Image as PILImage
import math
import sys
from pywinauto.application import Application
import win32gui
import win32con
import time
from win32api import GetSystemMetrics
from pptx import Presentation
from pptx.util import Inches
import os
from pptx.dml.color import RGBColor
from pptx.util import Pt

# instantiate an MCP server client
mcp = FastMCP("Calculator")

# DEFINE TOOLS

#addition tool
@mcp.tool()
def add(a: int, b: int) -> int:
    """Add two numbers"""
    print("CALLED: add(a: int, b: int) -> int:")
    return int(a + b)

@mcp.tool()
def add_list(l: list) -> int:
    """Add all numbers in a list"""
    print("CALLED: add(l: list) -> int:")
    return sum(l)

# subtraction tool
@mcp.tool()
def subtract(a: int, b: int) -> int:
    """Subtract two numbers"""
    print("CALLED: subtract(a: int, b: int) -> int:")
    return int(a - b)

# multiplication tool
@mcp.tool()
def multiply(a: int, b: int) -> int:
    """Multiply two numbers"""
    print("CALLED: multiply(a: int, b: int) -> int:")
    return int(a * b)

#  division tool
@mcp.tool() 
def divide(a: int, b: int) -> float:
    """Divide two numbers"""
    print("CALLED: divide(a: int, b: int) -> float:")
    return float(a / b)

# power tool
@mcp.tool()
def power(a: int, b: int) -> int:
    """Power of two numbers"""
    print("CALLED: power(a: int, b: int) -> int:")
    return int(a ** b)

# square root tool
@mcp.tool()
def sqrt(a: int) -> float:
    """Square root of a number"""
    print("CALLED: sqrt(a: int) -> float:")
    return float(a ** 0.5)

# cube root tool
@mcp.tool()
def cbrt(a: int) -> float:
    """Cube root of a number"""
    print("CALLED: cbrt(a: int) -> float:")
    return float(a ** (1/3))

# factorial tool
@mcp.tool()
def factorial(a: int) -> int:
    """factorial of a number"""
    print("CALLED: factorial(a: int) -> int:")
    return int(math.factorial(a))

# log tool
@mcp.tool()
def log(a: int) -> float:
    """log of a number"""
    print("CALLED: log(a: int) -> float:")
    return float(math.log(a))

# remainder tool
@mcp.tool()
def remainder(a: int, b: int) -> int:
    """remainder of two numbers divison"""
    print("CALLED: remainder(a: int, b: int) -> int:")
    return int(a % b)

# sin tool
@mcp.tool()
def sin(a: int) -> float:
    """sin of a number"""
    print("CALLED: sin(a: int) -> float:")
    return float(math.sin(a))

# cos tool
@mcp.tool()
def cos(a: int) -> float:
    """cos of a number"""
    print("CALLED: cos(a: int) -> float:")
    return float(math.cos(a))

# tan tool
@mcp.tool()
def tan(a: int) -> float:
    """tan of a number"""
    print("CALLED: tan(a: int) -> float:")
    return float(math.tan(a))

# mine tool
@mcp.tool()
def mine(a: int, b: int) -> int:
    """special mining tool"""
    print("CALLED: mine(a: int, b: int) -> int:")
    return int(a - b - b)

@mcp.tool()
def create_thumbnail(image_path: str) -> Image:
    """Create a thumbnail from an image"""
    print("CALLED: create_thumbnail(image_path: str) -> Image:")
    img = PILImage.open(image_path)
    img.thumbnail((100, 100))
    return Image(data=img.tobytes(), format="png")

@mcp.tool()
def strings_to_chars_to_int(string: str) -> list[int]:
    """Return the ASCII values of the characters in a word"""
    print("CALLED: strings_to_chars_to_int(string: str) -> list[int]:")
    return [int(ord(char)) for char in string]

@mcp.tool()
def int_list_to_exponential_sum(int_list: list) -> float:
    """Return sum of exponentials of numbers in a list"""
    print("CALLED: int_list_to_exponential_sum(int_list: list) -> float:")
    return sum(math.exp(i) for i in int_list)

@mcp.tool()
def fibonacci_numbers(n: int) -> list:
    """Return the first n Fibonacci Numbers"""
    print("CALLED: fibonacci_numbers(n: int) -> list:")
    if n <= 0:
        return []
    fib_sequence = [0, 1]
    for _ in range(2, n):
        fib_sequence.append(fib_sequence[-1] + fib_sequence[-2])
    return fib_sequence[:n]

@mcp.tool()
async def close_powerpoint() -> dict:
    """Close PowerPoint"""
    try:
        # Close PowerPoint
        os.system('taskkill /F /IM POWERPNT.EXE')
        time.sleep(2)
        
        return {
            "content": [
                TextContent(
                    type="text",
                    text="PowerPoint closed successfully"
                )
            ]
        }
    except Exception as e:
        print(f"Error in close_powerpoint: {str(e)}")
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error closing PowerPoint: {str(e)}"
                )
            ]
        }

@mcp.tool()
async def open_powerpoint() -> dict:
    """Open a new PowerPoint presentation"""
    try:
        # Close any existing PowerPoint instances
        await close_powerpoint()
        time.sleep(3)  # Increased wait time
        
        # Create a new presentation
        prs = Presentation()
        
        # Add a title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Save the presentation
        filename = 'presentation.pptx'
        prs.save(filename)
        time.sleep(5)  # Increased wait time for file save
        
        # Open the presentation
        os.startfile(filename)
        time.sleep(10)  # Increased wait time for PowerPoint to open
        
        return {
            "content": [
                TextContent(
                    type="text",
                    text="PowerPoint opened successfully with a new presentation"
                )
            ]
        }
    except Exception as e:
        print(f"Error in open_powerpoint: {str(e)}")
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error opening PowerPoint: {str(e)}"
                )
            ]
        }

@mcp.tool()
async def draw_rectangle(x1: int, y1: int, x2: int, y2: int) -> dict:
    """Draw a rectangle in the first slide of PowerPoint"""
    try:
        print(f"DEBUG: Drawing rectangle with raw parameters: x1={x1} ({type(x1)}), y1={y1} ({type(y1)}), x2={x2} ({type(x2)}), y2={y2} ({type(y2)})")
        
        # Convert parameters to integers
        try:
            x1 = int(float(str(x1)))
            y1 = int(float(str(y1)))
            x2 = int(float(str(x2)))
            y2 = int(float(str(y2)))
        except (ValueError, TypeError) as e:
            error_msg = f"Failed to convert parameters to integers: {str(e)}"
            print(f"DEBUG: {error_msg}")
            return {"content": [TextContent(type="text", text=error_msg)]}

        print(f"DEBUG: Converted coordinates: ({x1},{y1}) to ({x2},{y2})")
        
        # Validate coordinates
        if not (1 <= x1 <= 8 and 1 <= y1 <= 8 and 1 <= x2 <= 8 and 1 <= y2 <= 8):
            error_msg = f"Coordinates must be between 1 and 8, got: ({x1},{y1}) to ({x2},{y2})"
            print(f"DEBUG: {error_msg}")
            return {"content": [TextContent(type="text", text=error_msg)]}
        
        if x2 <= x1 or y2 <= y1:
            error_msg = f"End coordinates must be greater than start coordinates: ({x1},{y1}) to ({x2},{y2})"
            print(f"DEBUG: {error_msg}")
            return {"content": [TextContent(type="text", text=error_msg)]}
        
        # Wait before modifying the presentation
        time.sleep(2)
        
        # Ensure PowerPoint is closed before modifying the file
        await close_powerpoint()
        time.sleep(2)
        
        try:
            # Open the existing presentation
            prs = Presentation('presentation.pptx')
            slide = prs.slides[0]
            
            # Store existing text boxes
            text_boxes = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    text_boxes.append((text, left, top, width, height))
            
            # Clear existing shapes except text boxes
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    sp = shape._element
                    sp.getparent().remove(sp)
            
            # Convert coordinates to inches
            left = Inches(x1)
            top = Inches(y1)
            width = Inches(x2 - x1)
            height = Inches(y2 - y1)
            
            print(f"DEBUG: Rectangle dimensions - left={left}, top={top}, width={width}, height={height}")
            
            # Add rectangle
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE.RECTANGLE
                left, top, width, height
            )
            
            # Make the rectangle more visible
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White fill
            shape.line.color.rgb = RGBColor(0, 0, 0)  # Black border
            shape.line.width = Pt(4)  # Thicker border
            
            # Save the presentation
            prs.save('presentation.pptx')
            time.sleep(2)
            
            # Reopen PowerPoint
            os.startfile('presentation.pptx')
            time.sleep(5)
            
            print("DEBUG: Rectangle drawn successfully")
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=f"Rectangle drawn successfully from ({x1},{y1}) to ({x2},{y2})"
                    )
                ]
            }
            
        except Exception as e:
            error_msg = f"PowerPoint operation failed: {str(e)}"
            print(f"DEBUG: {error_msg}")
            return {"content": [TextContent(type="text", text=error_msg)]}
            
    except Exception as e:
        error_msg = f"Error in draw_rectangle: {str(e)}"
        print(f"DEBUG: {error_msg}")
        print(f"DEBUG: Error type: {type(e)}")
        import traceback
        traceback.print_exc()
        return {"content": [TextContent(type="text", text=error_msg)]}

@mcp.tool()
async def add_text_in_powerpoint(text: str) -> dict:
    """Add text to the first slide of PowerPoint"""
    try:
        print(f"DEBUG: Received text to add: {text}")
        print(f"DEBUG: Text type: {type(text)}")
        print(f"DEBUG: Text length: {len(text)}")
        print(f"DEBUG: Text contains newlines: {'\\n' in text}")
        
        # Wait before adding text
        time.sleep(5)
        
        # Ensure PowerPoint is closed before modifying the file
        await close_powerpoint()
        time.sleep(5)
        
        # Open the existing presentation
        prs = Presentation('presentation.pptx')
        slide = prs.slides[0]
        
        # Add a text box positioned inside the rectangle
        # Match the rectangle position from draw_rectangle
        left = Inches(2.2)  # Slightly more than rectangle left for margin
        top = Inches(2.5)   # Centered vertically in rectangle
        width = Inches(4.6) # Slightly less than rectangle width for margin
        height = Inches(2)  # Enough height for text
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()  # Clear existing text
        text_frame.word_wrap = True  # Enable word wrap
        text_frame.vertical_anchor = 1  # Middle vertical alignment
        
        # Split text into lines
        lines = text.split('\n')
        print(f"DEBUG: Number of lines: {len(lines)}")
        print(f"DEBUG: Lines to add: {lines}")
        
        # Add each line as a separate paragraph
        for i, line in enumerate(lines):
            if line.strip():  # Only add non-empty lines
                p = text_frame.add_paragraph()
                p.text = line.strip()
                p.alignment = 1  # Center align the text
                
                # Format the text
                run = p.runs[0]
                if "Final Result:" in line:
                    run.font.size = Pt(32)  # Header size
                    run.font.bold = True
                else:
                    run.font.size = Pt(28)  # Value size
                    run.font.bold = True
                
                run.font.color.rgb = RGBColor(0, 0, 0)  # Black text
                p.space_after = Pt(12)  # Add spacing between lines
        
        # Save and wait
        prs.save('presentation.pptx')
        time.sleep(5)
        
        # Reopen PowerPoint
        os.startfile('presentation.pptx')
        time.sleep(10)
        
        print(f"DEBUG: Text added successfully: {text}")
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Text added successfully: {text}"
                )
            ]
        }
    except Exception as e:
        print(f"Error in add_text_in_powerpoint: {str(e)}")
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error adding text: {str(e)}"
                )
            ]
        }

# DEFINE RESOURCES

# Add a dynamic greeting resource
@mcp.resource("greeting://{name}")
def get_greeting(name: str) -> str:
    """Get a personalized greeting"""
    print("CALLED: get_greeting(name: str) -> str:")
    return f"Hello, {name}!"


# DEFINE AVAILABLE PROMPTS
@mcp.prompt()
def review_code(code: str) -> str:
    return f"Please review this code:\n\n{code}"
    print("CALLED: review_code(code: str) -> str:")


@mcp.prompt()
def debug_error(error: str) -> list[base.Message]:
    return [
        base.UserMessage("I'm seeing this error:"),
        base.UserMessage(error),
        base.AssistantMessage("I'll help debug that. What have you tried so far?"),
    ]

if __name__ == "__main__":
    # Check if running with mcp dev command
    print("STARTING THE SERVER")
    if len(sys.argv) > 1 and sys.argv[1] == "dev":
        mcp.run()  # Run without transport for dev server
    else:
        mcp.run(transport="stdio")  # Run with stdio for direct execution
